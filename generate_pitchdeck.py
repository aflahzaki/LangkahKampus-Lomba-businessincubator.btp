#!/usr/bin/env python3
"""
Generate the complete PPTX pitch deck for LangkahKampus business case competition.
Populates 'Pitchdeck Business Idea Competition.pptx' with compelling pitch deck content.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Paths
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PPTX_FILE = os.path.join(SCRIPT_DIR, "Pitchdeck Business Idea Competition.pptx")
LOGO_FILE = os.path.join(SCRIPT_DIR, "Logo LangkahKampus.png")
TITLE_COLOR = RGBColor(0x2E, 0x40, 0x57)  # Dark blue #2E4057
BODY_COLOR = RGBColor(0x33, 0x33, 0x33)  # Dark gray for body
ACCENT_COLOR = RGBColor(0x1A, 0x73, 0xE8)  # Blue accent
STAT_COLOR = RGBColor(0xC0, 0x39, 0x2B)  # Red for key stats

# Slide dimensions (13.3" x 7.5")
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_title_textbox(slide, text, left=Inches(0.7), top=Inches(0.4),
                      width=Inches(12), height=Inches(1.0), font_size=Pt(32)):
    """Add a styled title text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.LEFT
    return txBox


def add_content_textbox(slide, left, top, width, height):
    """Add a content text box and return the text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_bullet_point(tf, text, level=0, font_size=Pt(18), bold=False,
                     color=None, space_before=Pt(6)):
    """Add a bullet point paragraph to text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color if color else BODY_COLOR
    p.space_before = space_before
    return p


def set_first_paragraph(tf, text, font_size=Pt(18), bold=False, color=None):
    """Set the first paragraph in a text frame."""
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color if color else BODY_COLOR
    return p


def clear_slide(slide):
    """Remove all shapes from a slide (except background)."""
    shapes_to_remove = list(slide.shapes)
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)


def build_slide_1(slide):
    """Title slide with logo and team info."""
    clear_slide(slide)

    # Title
    add_title_textbox(slide, "LangkahKampus",
                      left=Inches(0.7), top=Inches(0.5),
                      width=Inches(12), height=Inches(1.2),
                      font_size=Pt(44))

    # Subtitle
    tf = add_content_textbox(slide, Inches(0.7), Inches(1.6), Inches(12), Inches(0.8))
    set_first_paragraph(tf, "AI-Powered SNBP University Admission Advisory Platform",
                        font_size=Pt(24), color=RGBColor(0x55, 0x55, 0x55))

    # Logo - centered
    logo_path = LOGO_FILE
    if os.path.exists(logo_path):
        # Logo is 2816x1536, aspect ratio ~1.83:1
        logo_width = Inches(2.5)
        logo_height = Inches(1.36)  # maintain aspect ratio
        logo_left = (SLIDE_WIDTH - logo_width) // 2
        logo_top = Inches(2.8)
        slide.shapes.add_picture(logo_path, logo_left, logo_top,
                                 logo_width, logo_height)

    # Team info
    tf = add_content_textbox(slide, Inches(0.7), Inches(5.5), Inches(12), Inches(1.0))
    set_first_paragraph(tf, "Tim LangkahKampus | Telkom University",
                        font_size=Pt(20), color=RGBColor(0x55, 0x55, 0x55))
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Competition context
    tf2 = add_content_textbox(slide, Inches(0.7), Inches(6.3), Inches(12), Inches(0.6))
    set_first_paragraph(tf2, "Business Idea Competition - Bandung Techno Park 2025",
                        font_size=Pt(16), color=RGBColor(0x77, 0x77, 0x77))
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER


def build_slide_2(slide):
    """Masalah yang Diangkat."""
    clear_slide(slide)

    add_title_textbox(slide, "Masalah yang Diangkat",
                      left=Inches(0.7), top=Inches(0.4),
                      width=Inches(12), height=Inches(1.0),
                      font_size=Pt(32))

    # Key stat highlight
    tf_stat = add_content_textbox(slide, Inches(0.7), Inches(1.3), Inches(12), Inches(0.7))
    set_first_paragraph(tf_stat,
                        "800K pendaftar SNBP/tahun, hanya 27% diterima",
                        font_size=Pt(22), bold=True, color=STAT_COLOR)

    # Bullet points
    tf = add_content_textbox(slide, Inches(0.7), Inches(2.2), Inches(11.5), Inches(5.0))
    set_first_paragraph(tf,
                        "Choice-2 Trap: Universitas top menolak otomatis siswa yang "
                        "menempatkan mereka sebagai pilihan kedua",
                        font_size=Pt(20), bold=False, color=BODY_COLOR)

    add_bullet_point(tf,
                     "Data Asymmetry: Siswa tidak memiliki akses data historis acceptance rate "
                     "dan informasi kuota per program studi",
                     font_size=Pt(20), space_before=Pt(14))

    add_bullet_point(tf,
                     "Guru BK Overloaded: Rasio 1:300-500 siswa tanpa analytical tools "
                     "(70%+ tanpa digital tools)",
                     font_size=Pt(20), space_before=Pt(14))

    add_bullet_point(tf,
                     "Intra-School Collision: Siswa satu sekolah saling berkompetisi "
                     "tanpa sadar untuk slot universitas yang sama",
                     font_size=Pt(20), space_before=Pt(14))


def build_slide_3(slide):
    """Solusi yang Diberikan."""
    clear_slide(slide)

    add_title_textbox(slide, "Solusi: LangkahKampus",
                      left=Inches(0.7), top=Inches(0.4),
                      width=Inches(12), height=Inches(1.0),
                      font_size=Pt(32))

    # 5 Solution Pillars
    tf = add_content_textbox(slide, Inches(0.7), Inches(1.6), Inches(11.5), Inches(5.5))
    set_first_paragraph(tf, "5 Pilar Solusi:", font_size=Pt(20), bold=True,
                        color=ACCENT_COLOR)

    pillars = [
        "1. Hard-Rule Validator: Deteksi otomatis batasan Choice-2 universitas",
        "2. AI Probability Scoring: Ensemble XGBoost + LightGBM (AUC 0.91)",
        "3. What-If Engine: Counterfactual recommendations (DiCE + SHAP)",
        "4. Geospatial Profiling: Filter berdasarkan lokasi dan minat",
        "5. Anti-Bentrok Dashboard: De-confliction real-time untuk Guru BK",
    ]

    for pillar in pillars:
        add_bullet_point(tf, pillar, font_size=Pt(20), space_before=Pt(14))

    # Impact statement
    add_bullet_point(tf, "", font_size=Pt(10), space_before=Pt(20))
    add_bullet_point(tf,
                     "Dari data mentah menjadi rekomendasi actionable - "
                     "bukan hanya prediksi, tapi advisory lengkap",
                     font_size=Pt(18), bold=True, color=ACCENT_COLOR,
                     space_before=Pt(8))


def build_slide_4(slide):
    """Target Market (TAM, SAM, SOM)."""
    clear_slide(slide)

    add_title_textbox(slide, "Target Market",
                      left=Inches(0.7), top=Inches(0.4),
                      width=Inches(12), height=Inches(1.0),
                      font_size=Pt(32))

    # TAM
    tf = add_content_textbox(slide, Inches(0.7), Inches(1.6), Inches(11.5), Inches(5.5))
    set_first_paragraph(tf, "TAM (Total Addressable Market)",
                        font_size=Pt(20), bold=True, color=ACCENT_COLOR)
    add_bullet_point(tf, "5.9 juta siswa SMA/SMK | Rp154 miliar/tahun",
                     font_size=Pt(20), space_before=Pt(6))

    # SAM
    add_bullet_point(tf, "", font_size=Pt(8), space_before=Pt(14))
    add_bullet_point(tf, "SAM (Serviceable Addressable Market)",
                     font_size=Pt(20), bold=True, color=ACCENT_COLOR,
                     space_before=Pt(4))
    add_bullet_point(tf, "800K pendaftar SNBP | Rp57-65 miliar/tahun",
                     font_size=Pt(20), space_before=Pt(6))

    # SOM
    add_bullet_point(tf, "", font_size=Pt(8), space_before=Pt(14))
    add_bullet_point(tf, "SOM (Serviceable Obtainable Market)",
                     font_size=Pt(20), bold=True, color=ACCENT_COLOR,
                     space_before=Pt(4))
    add_bullet_point(tf, "Year 1: 50K siswa + 200 sekolah = Rp1.6 miliar",
                     font_size=Pt(20), space_before=Pt(6))
    add_bullet_point(tf, "Year 3: 200K siswa + 2000 sekolah = Rp10 miliar",
                     font_size=Pt(20), space_before=Pt(6))


def build_slide_5(slide):
    """Business Model."""
    clear_slide(slide)

    add_title_textbox(slide, "Business Model",
                      left=Inches(0.7), top=Inches(0.4),
                      width=Inches(12), height=Inches(1.0),
                      font_size=Pt(32))

    # Left column - B2C
    tf_left = add_content_textbox(slide, Inches(0.7), Inches(1.6),
                                  Inches(5.5), Inches(5.0))
    set_first_paragraph(tf_left, "B2C Freemium",
                        font_size=Pt(22), bold=True, color=ACCENT_COLOR)
    add_bullet_point(tf_left, "Free: Basic SNBP eligibility check",
                     font_size=Pt(18), space_before=Pt(10))
    add_bullet_point(tf_left, "Premium: Rp15-25K/prediksi",
                     font_size=Pt(18), space_before=Pt(8))
    add_bullet_point(tf_left, "  - Deep prediction + alternatives",
                     font_size=Pt(16), space_before=Pt(4))
    add_bullet_point(tf_left, "  - What-If counterfactual analysis",
                     font_size=Pt(16), space_before=Pt(4))

    # Right column - B2B
    tf_right = add_content_textbox(slide, Inches(6.8), Inches(1.6),
                                   Inches(5.8), Inches(5.0))
    set_first_paragraph(tf_right, "B2B SaaS - BK Command Center",
                        font_size=Pt(22), bold=True, color=ACCENT_COLOR)
    add_bullet_point(tf_right, "Basic: Rp2 juta/tahun (200 siswa)",
                     font_size=Pt(18), space_before=Pt(10))
    add_bullet_point(tf_right, "Pro: Rp3.5 juta/tahun (unlimited + analytics)",
                     font_size=Pt(18), space_before=Pt(8))
    add_bullet_point(tf_right, "Enterprise: Rp5 juta/tahun (API + dedicated)",
                     font_size=Pt(18), space_before=Pt(8))

    # Unit Economics at bottom
    tf_bottom = add_content_textbox(slide, Inches(0.7), Inches(5.8),
                                    Inches(12), Inches(1.2))
    set_first_paragraph(tf_bottom,
                        "Unit Economics: LTV:CAC ratio 8-12x | Gross Margin >85%",
                        font_size=Pt(20), bold=True, color=TITLE_COLOR)


def build_slide_6(slide):
    """Kompetitor dan Positioning."""
    clear_slide(slide)

    add_title_textbox(slide, "Kompetitor & Positioning",
                      left=Inches(0.7), top=Inches(0.4),
                      width=Inches(12), height=Inches(1.0),
                      font_size=Pt(32))

    # Comparison header
    tf = add_content_textbox(slide, Inches(0.7), Inches(1.5), Inches(12), Inches(0.5))
    set_first_paragraph(tf,
                        "LangkahKampus vs Ruangguru vs Aku Pintar vs Portal SNBP",
                        font_size=Pt(18), bold=True, color=BODY_COLOR)

    # Unique features - what only we have
    tf2 = add_content_textbox(slide, Inches(0.7), Inches(2.3), Inches(11.5), Inches(4.5))
    set_first_paragraph(tf2, "Fitur Unik LangkahKampus (yang tidak dimiliki kompetitor):",
                        font_size=Pt(20), bold=True, color=ACCENT_COLOR)

    features = [
        "Anti-Bentrok Dashboard - ONLY US (de-confliction intra-sekolah)",
        "Counterfactual AI (What-If Engine) - ONLY US",
        "Choice-2 Trap Detection - ONLY US",
        "SNBP-focused advisory: Rp15-25K vs kompetitor Rp100K+/bulan",
    ]
    for feat in features:
        add_bullet_point(tf2, feat, font_size=Pt(19), space_before=Pt(12))

    # Positioning statement
    add_bullet_point(tf2, "", font_size=Pt(8), space_before=Pt(16))
    add_bullet_point(tf2,
                     "Positioning: Advisory System (bukan sekadar data portal atau bimbel online)",
                     font_size=Pt(18), bold=True, color=TITLE_COLOR,
                     space_before=Pt(8))


def build_slide_7(slide):
    """Teknologi yang Digunakan."""
    clear_slide(slide)

    add_title_textbox(slide, "Teknologi",
                      left=Inches(0.7), top=Inches(0.4),
                      width=Inches(12), height=Inches(1.0),
                      font_size=Pt(32))

    tf = add_content_textbox(slide, Inches(0.7), Inches(1.5), Inches(11.5), Inches(5.5))

    set_first_paragraph(tf, "Tech Stack: Next.js | Node.js | FastAPI | PostgreSQL | Redis",
                        font_size=Pt(20), bold=True, color=ACCENT_COLOR)

    tech_items = [
        "ML: XGBoost + LightGBM ensemble (AUC 0.91, <5ms inference)",
        "XAI: SHAP + LIME + DiCE for explainable recommendations",
        "Infra: Docker + Kubernetes + GitHub Actions CI/CD",
        "Data: 47 engineered features, weekly refresh pipeline",
        "Security: JWT auth, rate limiting, data encryption at rest",
    ]

    for item in tech_items:
        add_bullet_point(tf, item, font_size=Pt(20), space_before=Pt(14))

    # Architecture flow
    add_bullet_point(tf, "", font_size=Pt(8), space_before=Pt(16))
    add_bullet_point(tf,
                     "Flow: User Input -> API Gateway -> ML Service -> "
                     "XAI Explanation -> Personalized Advisory",
                     font_size=Pt(18), bold=True, color=TITLE_COLOR,
                     space_before=Pt(8))


def build_slide_8(slide):
    """Analisis Keunggulan."""
    clear_slide(slide)

    add_title_textbox(slide, "Keunggulan Kompetitif",
                      left=Inches(0.7), top=Inches(0.4),
                      width=Inches(12), height=Inches(1.0),
                      font_size=Pt(32))

    tf = add_content_textbox(slide, Inches(0.7), Inches(1.5), Inches(11.5), Inches(5.5))

    set_first_paragraph(tf, "4 Keunggulan Utama:",
                        font_size=Pt(20), bold=True, color=ACCENT_COLOR)

    advantages = [
        "1. First-mover: Anti-Bentrok (tidak ada kompetitor yang memiliki fitur ini)",
        "2. Data Moat: Proprietary historical data yang terus bertambah setiap siklus",
        "3. Network Effect: Semakin banyak user per sekolah = prediksi semakin akurat",
        "4. Advisory vs Raw Data: Melengkapi (bukan bersaing dengan) portal pemerintah",
    ]

    for adv in advantages:
        add_bullet_point(tf, adv, font_size=Pt(20), space_before=Pt(14))

    # Defensive statement
    add_bullet_point(tf, "", font_size=Pt(8), space_before=Pt(16))
    add_bullet_point(tf,
                     "Defensive Position: Platform pemerintah menyediakan data mentah - "
                     "kami menyediakan intelligence dan rekomendasi actionable yang "
                     "tidak bisa dilakukan portal statis",
                     font_size=Pt(18), bold=True, color=TITLE_COLOR,
                     space_before=Pt(8))


def build_slide_9(slide):
    """Penutup / Closing."""
    clear_slide(slide)

    # Title
    add_title_textbox(slide, "Terima Kasih",
                      left=Inches(0.7), top=Inches(0.5),
                      width=Inches(12), height=Inches(1.2),
                      font_size=Pt(40))

    # Vision statement
    tf = add_content_textbox(slide, Inches(1.5), Inches(2.0), Inches(10.5), Inches(1.5))
    p = tf.paragraphs[0]
    p.text = ("\"Demokratisasi akses intelligence penerimaan universitas "
              "untuk setiap siswa Indonesia\"")
    p.font.size = Pt(22)
    p.font.italic = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Logo
    logo_path = LOGO_FILE
    if os.path.exists(logo_path):
        logo_width = Inches(2.5)
        logo_height = Inches(1.36)
        logo_left = (SLIDE_WIDTH - logo_width) // 2
        logo_top = Inches(3.8)
        slide.shapes.add_picture(logo_path, logo_left, logo_top,
                                 logo_width, logo_height)

    # Contact / next steps
    tf2 = add_content_textbox(slide, Inches(1.5), Inches(5.8), Inches(10.5), Inches(1.2))
    set_first_paragraph(tf2, "Tim LangkahKampus | Telkom University | Bandung",
                        font_size=Pt(18), color=RGBColor(0x55, 0x55, 0x55))
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_bullet_point(tf2, "Email: team@langkahkampus.id | Web: www.langkahkampus.id",
                     font_size=Pt(16), color=RGBColor(0x77, 0x77, 0x77),
                     space_before=Pt(8))
    tf2.paragraphs[-1].alignment = PP_ALIGN.CENTER


def build_slide_10(slide):
    """Slide 10: Referensi."""
    clear_slide(slide)

    add_title_textbox(slide, "Referensi",
                      left=Inches(0.7), top=Inches(0.3),
                      width=Inches(12), height=Inches(0.9),
                      font_size=Pt(32))

    # Left column - Data Sources
    tf_left = add_content_textbox(slide, Inches(0.7), Inches(1.3),
                                  Inches(6.0), Inches(5.8))
    set_first_paragraph(tf_left, "Sumber Data & Statistik:",
                        font_size=Pt(16), bold=True, color=ACCENT_COLOR)

    data_refs = [
        "[1] SNPMB/LTMPT (2024) - Statistik SNBP",
        "[2] BPS (2024) - Statistik Pendidikan Indonesia",
        "[3] Kemendikbudristek (2023) - Data Guru BK",
        "[4] DataReportal/We Are Social (2024) - Digital Indonesia",
        "[10] Google Trends (2024) - Tren 'Prediksi SNBP'",
        "[12] APJII (2024) - Penetrasi Internet Indonesia",
    ]

    for ref in data_refs:
        add_bullet_point(tf_left, ref, font_size=Pt(14), space_before=Pt(8))

    # Right column - Academic Papers
    tf_right = add_content_textbox(slide, Inches(6.9), Inches(1.3),
                                   Inches(6.0), Inches(5.8))
    set_first_paragraph(tf_right, "Publikasi Akademik:",
                        font_size=Pt(16), bold=True, color=ACCENT_COLOR)

    academic_refs = [
        "[5] Chen & Guestrin (2016) - XGBoost, KDD",
        "[6] Ke et al. (2017) - LightGBM, NeurIPS",
        "[7] Hollmann et al. (2023) - TabPFN, ICLR",
        "[8] Mothilal et al. (2020) - DiCE, FAT*",
        "[9] Lundberg & Lee (2017) - SHAP, NeurIPS",
        "[11] Ribeiro et al. (2016) - LIME, KDD",
    ]

    for ref in academic_refs:
        add_bullet_point(tf_right, ref, font_size=Pt(14), space_before=Pt(8))

    # Add reference image (ML model comparison) at the bottom
    img_path = os.path.join(SCRIPT_DIR, "references_images", "ml_model_comparison.png")
    if os.path.exists(img_path):
        img_width = Inches(4.5)
        img_left = (SLIDE_WIDTH - img_width) // 2
        slide.shapes.add_picture(img_path, img_left, Inches(5.5),
                                 img_width, Inches(1.8))


def main():
    """Main function to generate the pitch deck."""
    print("Loading template...")
    prs = Presentation(PPTX_FILE)

    slides = list(prs.slides)
    print(f"Template has {len(slides)} slides")

    # Build each slide
    builders = [
        build_slide_1,
        build_slide_2,
        build_slide_3,
        build_slide_4,
        build_slide_5,
        build_slide_6,
        build_slide_7,
        build_slide_8,
        build_slide_9,
        build_slide_10,
    ]

    # Add extra slides if needed
    while len(slides) < len(builders):
        slide_layout = prs.slide_layouts[5]  # Blank layout
        new_slide = prs.slides.add_slide(slide_layout)
        slides = list(prs.slides)
        print(f"  Added new slide (total: {len(slides)})")

    if len(slides) < len(builders):
        raise ValueError(
            f"Template has {len(slides)} slides but {len(builders)} are required. "
            "Please ensure the PPTX template contains at least 10 slides."
        )

    for i, builder in enumerate(builders):
        print(f"Building slide {i + 1}...")
        builder(slides[i])

    # Save
    print(f"Saving to {PPTX_FILE}...")
    prs.save(PPTX_FILE)
    print("Done! Pitch deck generated successfully.")


if __name__ == "__main__":
    main()
